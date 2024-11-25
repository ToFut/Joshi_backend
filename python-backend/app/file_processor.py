import base64
import os
import tempfile
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
import mammoth
import pptx
import shutil

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        raise ValueError(f"Error extracting text from TXT file: {e}")

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise ValueError(f"Error extracting text from DOCX file: {e}")

def extract_text_from_pdf(file_path):
    try:
        return extract_pdf_text(file_path)
    except Exception as e:
        raise ValueError(f"Error extracting text from PDF file: {e}")

def extract_text_from_doc(file_path):
    try:
        with open(file_path, "rb") as doc_file:
            result = mammoth.extract_raw_text(doc_file)
            return result.value
    except Exception as e:
        raise ValueError(f"Error extracting text from DOC file: {e}")

def extract_text_from_pptx(file_path):
    try:
        presentation = pptx.Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Error extracting text from PPTX file: {e}")

def extract_text_from_file(file_name, temp_file_path):
    temp_dir = tempfile.mkdtemp()
    try:
        file_ext = file_name.split('.')[-1].lower()
        # Extract text based on file extension
        if file_ext == 'txt':
            extracted_text = extract_text_from_txt(temp_file_path)
        elif file_ext == 'docx':
            extracted_text = extract_text_from_docx(temp_file_path)
        elif file_ext == 'pdf':
            extracted_text = extract_text_from_pdf(temp_file_path)
        elif file_ext == 'doc':
            extracted_text = extract_text_from_doc(temp_file_path)
        elif file_ext == 'pptx':
            extracted_text = extract_text_from_pptx(temp_file_path)
        else:
            return {'success': False, 'error': 'Unsupported file type'}
        
        print(extracted_text[:100])
        return {'success': True, 'content': extracted_text}

    except Exception as e:
        return {'success': False, 'error': f"Error processing file: {e}"}
    finally:
        # Cleanup temp directory
        shutil.rmtree(temp_dir, ignore_errors=True)
